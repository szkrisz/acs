from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# --------- Style (Ice) ----------
ICE_BLUE = RGBColor(0xAD, 0xD8, 0xE6)   # #ADD8E6
ICE_BLUE_DARK = RGBColor(0x5A, 0xA9, 0xC9)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x00, 0x00, 0x00)
GRAY     = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF8)

FONT_NAME = "Calibri"
TITLE_SIZE = Pt(34)
SUBTITLE_SIZE = Pt(18)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(14)
FOOTER_SIZE = Pt(10)

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

def add_ice_header(slide, title_text, subtitle_text=None):
    # Header bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ICE_BLUE
    bar.line.color.rgb = ICE_BLUE

    # subtle bottom line
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.0), Inches(13.333), Inches(0.04)
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = ICE_BLUE_DARK
    ln.line.color.rgb = ICE_BLUE_DARK

    # Title
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12.2), Inches(0.6))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_NAME
    run.font.size = TITLE_SIZE
    run.font.bold = True
    run.font.color.rgb = BLACK

    if subtitle_text:
        st = slide.shapes.add_textbox(Inches(0.62), Inches(1.12), Inches(12.0), Inches(0.5))
        tf2 = st.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.name = FONT_NAME
        r2.font.size = SUBTITLE_SIZE
        r2.font.color.rgb = BLACK

def add_footer(slide, text):
    ft = slide.shapes.add_textbox(Inches(0.6), Inches(7.12), Inches(12.2), Inches(0.3))
    tf = ft.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = FONT_NAME
    r.font.size = FOOTER_SIZE
    r.font.color.rgb = GRAY

def add_bullets(slide, bullets, top=1.55, left=0.9, width=11.9, height=5.3, font_size=BODY_SIZE):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.name = FONT_NAME
        p.font.size = font_size
        p.font.color.rgb = BLACK
        p.level = 0
        p.space_after = Pt(6)

def add_note(slide, note_text):
    # Speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = note_text

def icon_box(slide, x, y, w, h, title, body):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.color.rgb = ICE_BLUE_DARK
    rect.line.width = Pt(1.5)

    # Title
    t = slide.shapes.add_textbox(Inches(x+0.25), Inches(y+0.18), Inches(w-0.5), Inches(0.35))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT_NAME
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = BLACK

    # Body
    b = slide.shapes.add_textbox(Inches(x+0.25), Inches(y+0.58), Inches(w-0.5), Inches(h-0.75))
    tfb = b.text_frame
    tfb.clear()
    p2 = tfb.paragraphs[0]
    p2.text = body
    p2.font.name = FONT_NAME
    p2.font.size = Pt(13)
    p2.font.color.rgb = BLACK
    p2.space_after = Pt(2)

def arrow(slide, x1, y1, x2, y2):
    # simple thin line + triangle to mimic arrow
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y1), Inches(x2-x1), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ICE_BLUE_DARK
    line.line.color.rgb = ICE_BLUE_DARK

    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x2-0.15), Inches(y1-0.08), Inches(0.25), Inches(0.2))
    tri.fill.solid()
    tri.fill.fore_color.rgb = ICE_BLUE_DARK
    tri.line.color.rgb = ICE_BLUE_DARK

def roadmap_timeline(slide, y=2.1):
    # 3 phase boxes
    icon_box(slide, 0.8, y, 3.9, 1.55, "0–6 hó", "MVP + protó\nPilot előkészítés\nSzerződés + spare + connectivity")
    icon_box(slide, 4.75, y, 3.9, 1.55, "6–12 hó", "Pilot (6–8)\nMarketing indul (8)\nSales rámpa (9)\nLOI/pipeline")
    icon_box(slide, 8.7, y, 3.85, 1.55, "12–24 hó", "2–6 fizetős telepítés\nKiállítás 2. év\nPartner bővítés\nEU előkészítés")
    # arrows
    arrow(slide, 4.15, y+0.78, 4.75, y+0.78)
    arrow(slide, 8.1, y+0.78, 8.7, y+0.78)

def pricing_cards(slide, y=2.1):
    icon_box(slide, 0.8, y, 4.1, 1.6, "Telepítés", "3–10M Ft / telepítés\nCommissioning benne\nTervezés külön")
    icon_box(slide, 5.0, y, 3.75, 1.6, "CLOUD", "150k Ft / site / hó\nLOCAL-ra épül\nMagasabb margin")
    icon_box(slide, 8.85, y, 3.7, 1.6, "3 értékesítési út", "Integráció\nUpgrade\nTurnkey")

def local_cloud_diagram(slide, y=2.0):
    icon_box(slide, 0.9, y, 5.9, 2.0, "LOCAL (helyben)", "Felügyelet + riport\nEseménynapló\nManuális beavatkozás támogatás\nAlacsony komplexitás")
    icon_box(slide, 7.05, y, 5.35, 2.0, "CLOUD (előfizetés)", "Energia menedzsment\nIdőjárási kontextus\nPrediktív előrejelzés\nEmelt support (2–5 perc)")
    arrow(slide, 6.6, y+1.02, 7.05, y+1.02)

def three_modes_diagram(slide, y=2.1):
    icon_box(slide, 0.8, y, 4.1, 1.75, "A) Integráció", "Meglévő rendszerbe\nMinimál HW\nGyors bevezetés")
    icon_box(slide, 4.75, y, 4.1, 1.75, "B) Upgrade", "HW+SW frissítés\nLegjobb belépő érték\nLOCAL→CLOUD")
    icon_box(slide, 8.7, y, 3.85, 1.75, "C) Turnkey", "Új hűtőház\nTeljes csomag\nPartner kivitelezés")

def kpi_gauges(slide, y=2.15):
    # simple "badge" style KPIs
    def badge(x, title, value, desc):
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.95), Inches(1.45))
        b.fill.solid()
        b.fill.fore_color.rgb = WHITE
        b.line.color.rgb = ICE_BLUE_DARK
        b.line.width = Pt(2)

        t = slide.shapes.add_textbox(Inches(x+0.25), Inches(y+0.18), Inches(3.4), Inches(0.35))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = BLACK

        v = slide.shapes.add_textbox(Inches(x+0.25), Inches(y+0.55), Inches(3.6), Inches(0.5))
        tf2 = v.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = value
        r2.font.name = FONT_NAME
        r2.font.size = Pt(28)
        r2.font.bold = True
        r2.font.color.rgb = ICE_BLUE_DARK

        d = slide.shapes.add_textbox(Inches(x+0.25), Inches(y+1.05), Inches(3.6), Inches(0.35))
        tf3 = d.text_frame
        tf3.clear()
        p3 = tf3.paragraphs[0]
        p3.text = desc
        p3.font.name = FONT_NAME
        p3.font.size = Pt(11)
        p3.font.color.rgb = GRAY

    badge(0.85, "Pilot uptime", "100%", "30 nap: nincs adatkimaradás és a vezérlés sem áll le")
    badge(4.75, "Reakcióidő", "2–5 perc", "Riasztás → első intézkedés / diagnózis megkezdése")
    badge(8.65, "Deliverables", "1 case", "Anonimizált esettanulmány + checklist + üzemeltetési v1")

def make_slide(prs, title, bullets=None, subtitle=None, visual=None, notes=""):
    s = blank_slide(prs)
    add_ice_header(s, title, subtitle)
    if visual:
        visual(s)
    if bullets:
        add_bullets(s, bullets, top=4.05, left=0.9, width=11.9, height=2.6, font_size=SMALL_SIZE)
    add_footer(s, prs._footer_text)
    add_note(s, notes)
    return s

def make_text_slide(prs, title, bullets, subtitle=None, notes=""):
    s = blank_slide(prs)
    add_ice_header(s, title, subtitle)
    add_bullets(s, bullets, top=1.65, left=0.95, width=12.0, height=5.2, font_size=BODY_SIZE)
    add_footer(s, prs._footer_text)
    add_note(s, notes)
    return s

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 wide
    prs.slide_height = Inches(7.5)
    prs._footer_text = "CUC – Cooling Under Control | Befektetői deck | 2026-03-30"

    # Slide 1 (Title)
    s = blank_slide(prs)
    add_ice_header(s, "CUC – Cooling Under Control", "Integrált cold storage control: LOCAL megbízhatóság + CLOUD intelligencia")
    # hero cards
    pricing_cards(s, y=2.25)
    add_bullets(s, [
        "Pilot: 6–8. hónap (Dombegyház) → marketing indul 8. hónaptól → skálázás HU, majd EU előkészítés",
        "Cél: stabil üzemeltetés, gyors reakció és átláthatóság (LOCAL), majd energia/predikció (CLOUD)",
    ], top=4.35, font_size=SMALL_SIZE)
    add_footer(s, prs._footer_text)
    add_note(s, "Nyitó: 20–30 mp. Mi a CUC és miért számít: a hűtőházi üzemeltetés kontrollja (adat+vezérlés) és skálázható cloud érték.")

    # Slide 2 (Problem)
    make_text_slide(
        prs,
        "A probléma",
        [
            "Sok szereplő (kivitelező, vezérlés, szoftver) → nincs egy felelős, nehéz hibát izolálni",
            "Rejtett hibák éjszaka/hétvégén → késői észlelés → romlás és kárkockázat",
            "Kevés adat és visszakereshetőség → energia és minőség ingadozik",
            "Az ügyfél nem IT-projektet akar, hanem stabil tárolást és gyors reagálást"
        ],
        notes="A fájdalom: felelősség hiánya + késői észlelés. A vevő a kockázatot viseli (áru, energia), de nincs eszköze kontrollálni."
    )

    # Slide 3 (Solution)
    make_text_slide(
        prs,
        "A megoldás: CUC platform",
        [
            "Vezérlés + felügyelet + dashboard egy rendszerben",
            "Riasztások, eseménynapló, távoli diagnózis és beavatkozás (igény szerint)",
            "Workflow szemlélet (termény-specifikus logika később bővíthető)",
            "Minősített komponensekkel indulunk → gyors time-to-market, alacsonyabb megfelelőségi kockázat"
        ],
        notes="Egy felelős réteg az üzemeltetés felett: monitorozás + kontroll. Nem dobozos HW-termékkel indulunk, hanem gyors integrációval."
    )

    # Slide 4 (LOCAL vs CLOUD) with diagram
    make_slide(
        prs,
        "LOCAL vs CLOUD – két lépcsős termék",
        bullets=[
            "LOCAL: helyi felügyelet és riport; alacsony komplexitás",
            "CLOUD: energia menedzsment + időjárási kontextus + predikció + emelt support",
            "Ügyfél út: LOCAL belépés → CLOUD előfizetés bizonyított érték után"
        ],
        visual=lambda slide: local_cloud_diagram(slide, y=1.95),
        notes="LOCAL csökkenti a belépési kockázatot. CLOUD a skálázható, magas marginú réteg (150k/site/hó)."
    )

    # Slide 5 (Value/KPI) with badges
    make_slide(
        prs,
        "Értékajánlat (mit javítunk)",
        bullets=[
            "Reakcióidő cél: 2–5 perc (riasztás → diagnózis/első lépés)",
            "Auditálhatóság: eseménynapló és trendek – „mi történt és mikor?”",
            "CLOUD szint: energiaoptimalizálás és predikció (több site-ra skálázható)"
        ],
        visual=lambda slide: kpi_gauges(slide, y=2.05),
        notes="A riasztások száma nem KPI. A KPI: stabil üzem + gyors reakció + visszakereshetőség."
    )

    # Slide 6 (3 modes) with cards
    make_slide(
        prs,
        "3 értékesítési mód (go-to-market rugalmasság)",
        bullets=[
            "A) Integráció: meglévő rendszerbe, minimál HW, gyors bevezetés",
            "B) Upgrade: HW+SW frissítés, tipikusan legjobb belépő érték",
            "C) Turnkey: új hűtőház, teljes csomag partner kivitelezéssel"
        ],
        visual=lambda slide: three_modes_diagram(slide, y=2.05),
        notes="Nem csak greenfield: integráció/upgrade gyors piacnyitás, turnkey nagyobb kosárérték."
    )

    # Slide 7 (Pricing) with cards
    make_slide(
        prs,
        "Árazás és bevételi modell",
        bullets=[
            "Telepítés: 3–10M Ft / telepítés (scope és funkció szerint)",
            "Telepítés/commissioning benne van; tervezés (engineering/design) külön megállapodás szerint",
            "CLOUD: 150 000 Ft / site / hó (LOCAL-ra épül)",
            "Pozícionálás: jellemzően a teljes beruházás 10–15%-a"
        ],
        visual=lambda slide: pricing_cards(slide, y=2.05),
        notes="Projektár ad kezdeti cashflow-t, CLOUD ad ismétlődő bevételt (ARR) és magasabb marginú skálázást."
    )

    # Slide 8 (Roadmap timeline)
    make_slide(
        prs,
        "Roadmap – 0–6 / 6–12 / 12–24 hónap",
        bullets=[
            "0–6 hó: MVP + protó; pilot előkészítés (szerződés, spare, connectivity)",
            "6–12 hó: pilot (6–8) + 30 nap stabil üzem; marketing 8-tól; sales 9-től",
            "12–24 hó: 2–6 fizetős telepítés HU-n; kiállítás 2. év; partner bővítés; EU előkészítés"
        ],
        visual=lambda slide: roadmap_timeline(slide, y=2.05),
        notes="Időzítés lényege: előbb bizonyíték (pilot), utána költés (marketing/sales), kiállítás csak akkor, ha van social proof."
    )

    # Slide 9 (Pilot success) with KPI badges
    make_slide(
        prs,
        "Pilot: siker definíció (Dombegyház)",
        bullets=[
            "Uptime 30 nap: 100% – nincs adatkimaradás és a vezérlés sem áll le",
            "Support reakcióidő cél: 2–5 perc",
            "Deliverables: anonimizált case study + üzemeltetési v1 + telepítési checklist + SLA vázlat"
        ],
        visual=lambda slide: kpi_gauges(slide, y=2.05),
        notes="Pilot cél: üzembiztos referencia. 100% uptime itt adat+vezérlés értelmű. A puffer részben ezt védi (spare/RMA/travel)."
    )

    # Slide 10 (GTM) with simple funnel cards
    def gtm_visual(slide):
        icon_box(slide, 0.8, 2.0, 4.1, 1.6, "Proof", "Pilot case study\nReferenciák\nMérhető KPI")
        icon_box(slide, 4.75, 2.0, 4.1, 1.6, "Acquire", "Célzott hirdetés\nWeb + tartalom\nDirect outreach")
        icon_box(slide, 8.7, 2.0, 3.85, 1.6, "Convert", "Ajánlati csomag\nSLA vázlat\nPartner kivitelezés")
        arrow(slide, 4.15, 2.8, 4.75, 2.8)
        arrow(slide, 8.1, 2.8, 8.7, 2.8)

    make_slide(
        prs,
        "Go-to-market (8. hónaptól, kiállítás nélkül)",
        bullets=[
            "Marketing 8. hónaptól: lead generálás + sales enablement (pilot proof után)",
            "Partner kivitelezők csatornaként",
            "Kiállítás csak a 2. évben, amikor már van social proof"
        ],
        visual=gtm_visual,
        notes="A marketing ajánlat (12M+ÁFA/év) csak pilot után indul. Kiállítás: 2. év, amikor már van mit bemutatni."
    )

    # Slide 11 (Cost plan + buffer) with blocks
    def cost_visual(slide):
        icon_box(slide, 0.8, 2.0, 3.95, 1.7, "Fix run-rate", "CEO + mérnök\nMinimál rezsi/IT\nAdmin")
        icon_box(slide, 4.8, 2.0, 3.95, 1.7, "Pilot csúcs (6–8)", "HW + travel\nSpare kit (RMA)\nEgyszeri jogi")
        icon_box(slide, 8.8, 2.0, 3.75, 1.7, "Skálázás", "Marketing 8-tól\nSales 9-től\nKiállítás 2. év")
        # buffer ribbon
        rib = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.85), Inches(11.75), Inches(0.5))
        rib.fill.solid()
        rib.fill.fore_color.rgb = ICE_BLUE
        rib.line.color.rgb = ICE_BLUE_DARK
        tx = slide.shapes.add_textbox(Inches(1.05), Inches(3.92), Inches(11.2), Inches(0.35))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "A kért összeg pufferként szolgál a fix kiadások mellett (csúszás, RMA, cashflow, skálázási csúcsok)."
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLACK

    make_slide(
        prs,
        "Költségterv – fix run-rate + pilot csúcs + skálázás",
        bullets=[
            "Marketing ügynökség: ~12M Ft + ÁFA / év (8. hónaptól indul)",
            "Kiállítás: csak a 2. évben (+3–4.5M + ÁFA, traction után)"
        ],
        visual=cost_visual,
        notes="Költségnarratíva: lean fix költség, pilot csúcs, majd kontrollált skálázás. A puffer biztosítja, hogy a pilot cél (100% uptime) ne sérüljön."
    )

    # Slide 12 (Ask) with simple tranche boxes
    def ask_visual(slide):
        icon_box(slide, 0.8, 2.0, 5.8, 1.8, "Ask (seed)", "HU validáció + skálázás\nPilot → fizetős telepítések\nCLOUD ARR alapok")
        icon_box(slide, 6.75, 2.0, 5.8, 1.8, "Tranche (opcionális)", "T1: pilot + alap működés\nT2: pilot siker + pipeline\n→ skálázó költés")
        rib = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.05), Inches(11.75), Inches(1.1))
        rib.fill.solid()
        rib.fill.fore_color.rgb = LIGHT_GRAY
        rib.line.color.rgb = ICE_BLUE_DARK
        tx = slide.shapes.add_textbox(Inches(1.1), Inches(4.15), Inches(11.2), Inches(0.9))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Puffer célja: gyors iteráció és kockázatkezelés (RMA/csere, kiszállás, ellátási lánc, 30–60 napos fizetések)."
        p.font.name = FONT_NAME
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLACK

    make_slide(
        prs,
        "A befektetés (Ask) – és miért kell puffer",
        bullets=[
            "Pilot cél biztosítása: 30 nap 100% uptime (adat + vezérlés) és 2–5 perc reakcióidő",
            "A puffer nem kényelmi: a pilot és az első értékesítések nem állhatnak meg váratlan csúcsokon"
        ],
        visual=ask_visual,
        notes="Zárás: a tőkét fegyelmezetten költjük. A puffer az uptime cél és a gyors iteráció garanciája. Következő kör alapja: fizetős telepítések + ARR."
    )

    # Backup 1: Pricing summary with ARR example
    make_slide(
        prs,
        "Backup: Pricing & ARR példa",
        bullets=[
            "Telepítés: 3–10M Ft (commissioning benne; tervezés külön)",
            "CLOUD: 150k Ft / site / hó (LOCAL-ra épül)",
            "Példa: 10 site × 150k Ft/hó = 1.5M Ft/hó = 18M Ft/év (CLOUD bevétel)"
        ],
        visual=lambda slide: pricing_cards(slide, y=2.05),
        notes="Ha rákérdeznek: a CLOUD bevétel a skálázható rész. 10 site már érezhető ARR, és minden új site növeli a moatat (adat)."
    )

    # Backup 2: Q&A
    make_text_slide(
        prs,
        "Backup: Q&A – tipikus befektetői kérdések",
        [
            "Mi az első beachhead szegmens és miért fizet? (kárkockázat + energia + felelősség)",
            "Mi a bruttó margin logika telepítésnél és a CLOUD előfizetésnél?",
            "Mi a moat? (workflow + adatok + integráció + standard telepítési csomag)",
            "Mi a legnagyobb kockázat és mitigáció? (RMA/spare, szerződés, 2. kivitelező, security alapok)",
            "Mi a következő kör mérföldköve? (pilot + fizetős telepítések + pipeline + ARR)"
        ],
        notes="Tartsd röviden: 1 mondat válasz + 1 konkrét példa. Ha számot kérnek: pricing és 150k/site/hó már a deckben van."
    )

    out = "CUC_Pitch_Deck_ICE_HU.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()
