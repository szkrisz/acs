from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --------- Style (Ice) ----------
ICE_BLUE = RGBColor(0xAD, 0xD8, 0xE6)   # #ADD8E6
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x00, 0x00, 0x00)

FONT_NAME = "Calibri"
TITLE_SIZE = Pt(34)
SUBTITLE_SIZE = Pt(18)
BODY_SIZE = Pt(18)
FOOTER_SIZE = Pt(10)

def add_ice_header(slide, title_text, subtitle_text=None):
    """Add ice style header bar and title."""
    # Header bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ICE_BLUE
    bar.line.color.rgb = ICE_BLUE

    # Title
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12.2), Inches(0.6))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_NAME
    run.font.size = TITLE_SIZE
    run.font.bold = True
    run.font.color.rgb = BLACK

    # Subtitle (optional)
    if subtitle_text:
        st = slide.shapes.add_textbox(Inches(0.62), Inches(1.1), Inches(12.0), Inches(0.5))
        tf2 = st.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.name = FONT_NAME
        r2.font.size = SUBTITLE_SIZE
        r2.font.color.rgb = BLACK

def add_bullets(slide, bullets, top=1.55, left=0.9, width=11.9, height=5.3):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.name = FONT_NAME
        p.font.size = BODY_SIZE
        p.font.color.rgb = BLACK
        p.level = 0
        p.space_after = Pt(6)

def add_footer(slide, text):
    ft = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(12.2), Inches(0.3))
    tf = ft.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = FONT_NAME
    r.font.size = FOOTER_SIZE
    r.font.color.rgb = BLACK

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # White background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 wide
    prs.slide_height = Inches(7.5)

    footer = "CUC – Cooling Under Control | Befektetői deck (vázlat) | 2026-03-30"

    # --- Slide 1 ---
    s = blank_slide(prs)
    add_ice_header(s, "CUC – Cooling Under Control", "Integrált cold storage control: LOCAL megbízhatóság + CLOUD intelligencia")
    add_bullets(s, [
        "Cél: stabil üzemeltetés + gyors reakció + átláthatóság zöldség/hagyma tárolóknál",
        "Kétlépcsős termék: LOCAL (helyben) → CLOUD (energia, predikció, emelt support)",
        "Pilot: 6–8. hónap (Dombegyház) → marketing indul 8. hónaptól → skálázás HU, majd EU előkészítés"
    ], top=2.1)
    add_footer(s, footer)

    # --- Slide 2: Probléma ---
    s = blank_slide(prs)
    add_ice_header(s, "A probléma")
    add_bullets(s, [
        "Sok szereplő (kivitelező, vezérlés, szoftver) → nincs “single owner”, nehéz hibát izolálni",
        "Rejtett hibák éjszaka/hétvégén → késői észlelés → romlás és kárkockázat",
        "Kevés adat és visszakereshetőség → energia és minőség ingadozik",
        "Az ügyfél nem IT-projektet akar, hanem stabil tárolást"
    ])
    add_footer(s, footer)

    # --- Slide 3: Megoldás ---
    s = blank_slide(prs)
    add_ice_header(s, "A megoldás: CUC platform")
    add_bullets(s, [
        "Vezérlés + felügyelet + dashboard egy rendszerben",
        "Riasztások, eseménynapló, távoli diagnózis és beavatkozás (igény szerint)",
        "Workflow szemlélet (termény-specifikus logika később bővíthető)",
        "Minősített komponensekkel indulunk → gyors time-to-market"
    ])
    add_footer(s, footer)

    # --- Slide 4: LOCAL vs CLOUD ---
    s = blank_slide(prs)
    add_ice_header(s, "LOCAL vs CLOUD – két lépcsős termék")
    add_bullets(s, [
        "LOCAL: minden szolgáltatás helyben; felügyelet + riport; automatizált hibaelhárítás nélkül",
        "CLOUD: LOCAL + felhő; energia menedzsment; időjárás; prediktív előrejelzés; emelt support",
        "Ügyfél út: gyors belépés LOCAL-lal → bizonyított érték után CLOUD előfizetés"
    ])
    add_footer(s, footer)

    # --- Slide 5: Érték / KPI ---
    s = blank_slide(prs)
    add_ice_header(s, "Értékajánlat (mit javítunk)")
    add_bullets(s, [
        "Reakcióidő cél: 2–5 perc (riasztás → intézkedés/diagnózis)",
        "Auditálhatóság: eseménynapló és trendek – „mi történt és mikor?”",
        "Stabil üzemeltetés és kockázatcsökkentés",
        "CLOUD: energiaoptimalizálás és predikció (skálázható érték)"
    ])
    add_footer(s, footer)

    # --- Slide 6: 3 értékesítési mód ---
    s = blank_slide(prs)
    add_ice_header(s, "3 értékesítési mód")
    add_bullets(s, [
        "A) Integráció meglévő rendszerbe: minimál HW + LOCAL",
        "B) Upgrade: HW + SW frissítés",
        "C) Turnkey: új hűtőház, teljes csomag partner kivitelezéssel"
    ])
    add_footer(s, footer)

    # --- Slide 7: Árazás ---
    s = blank_slide(prs)
    add_ice_header(s, "Árazás és bevételi modell")
    add_bullets(s, [
        "Telepítés: 3–10M Ft / telepítés (scope szerint), commissioning benne",
        "Tervezés (engineering/design) külön megállapodás szerint",
        "CLOUD előfizetés: 150 000 Ft / site / hó",
        "Pozícionálás: jellemzően a teljes beruházás 10–15%-a"
    ])
    add_footer(s, footer)

    # --- Slide 8: Roadmap ---
    s = blank_slide(prs)
    add_ice_header(s, "Roadmap – 0–6 / 6–12 / 12–24 hónap")
    add_bullets(s, [
        "0–6 hó: MVP + protó; pilot előkészítés",
        "6–12 hó: pilot (6–8) + 30 nap stabil üzem; marketing 8-tól; sales 9-től",
        "12–24 hó: 2–6 fizetős telepítés HU-n; kiállítás 2. év; partner bővítés; EU előkészítés"
    ])
    add_footer(s, footer)

    # --- Slide 9: Pilot KPI ---
    s = blank_slide(prs)
    add_ice_header(s, "Pilot: siker definíció")
    add_bullets(s, [
        "Uptime 30 nap: 100% – nincs adatkimaradás és a vezérlés sem áll le",
        "Support reakcióidő cél: 2–5 perc",
        "Deliverables: case study + üzemeltetési v1 + telepítési checklist + SLA vázlat"
    ])
    add_footer(s, footer)

    # --- Slide 10: GTM / Marketing ---
    s = blank_slide(prs)
    add_ice_header(s, "Go-to-market (8. hónaptól, kiállítás nélkül)")
    add_bullets(s, [
        "Pilot case study → bizalom és ajánlások",
        "Célzott online jelenlét + hirdetés + sales enablement anyagok",
        "Direct outreach + partner kivitelezők mint csatorna",
        "Kiállítás csak a 2. évben (social proof után)"
    ])
    add_footer(s, footer)

    # --- Slide 11: Költségterv + marketing ajánlat ---
    s = blank_slide(prs)
    add_ice_header(s, "Költségterv – fix run-rate + pilot csúcs + skálázás")
    add_bullets(s, [
        "Fix működés: core csapat + minimál OPEX",
        "Pilot csúcsok (6–8): hardver, travel/szállás, spare kit, egyszeri jogi",
        "Marketing ügynökség: ~12M Ft + ÁFA / év (8. hónaptól), kiállítás csak 2. év (+3–4.5M + ÁFA)",
        "A kért összeg pufferként szolgál a fix kiadások mellett (csúszás, RMA, cashflow, skálázási csúcsok)"
    ])
    add_footer(s, footer)

    # --- Slide 12: Ask + Buffer ---
    s = blank_slide(prs)
    add_ice_header(s, "Ask – és miért kell puffer")
    add_bullets(s, [
        "Seed tőke: HU validáció + skálázás finanszírozása (pilot → fizetős telepítések)",
        "Puffer célja: pilot 30 napos 100% uptime tartása (adat + vezérlés) és gyors iteráció",
        "Puffer lefedi: RMA/csere, kiszállások, ellátási lánc, 30–60 napos fizetések, marketing/sales ramp-up",
        "Tranche (opcionális): T1 pilot + alap működés; T2 pilot siker + pipeline után skálázó költés"
    ])
    add_footer(s, footer)

    # Backups omitted in v1
    out = "CUC_Pitch_Deck_ICE_HU.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()
